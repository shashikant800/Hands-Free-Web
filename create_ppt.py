#!/usr/bin/env python3
"""
Generate 8-slide PowerPoint presentation for Nutshell - HackJNU4.0 Hackathon Submission
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(139, 69, 19)      # Brown (nutshell color)
ACCENT = RGBColor(245, 166, 35)      # Orange/Gold
DARK_BG = RGBColor(30, 30, 30)       # Dark background
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)
DARK_TEXT = RGBColor(50, 50, 50)

def add_title_slide(prs, title, subtitle):
    """Slide 1: Title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BG
    bg_shape.line.fill.background()
    
    # Emoji icon
    icon_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(2.5), Inches(1.5))
    icon_tf = icon_box.text_frame
    icon_p = icon_tf.paragraphs[0]
    icon_p.text = "🥜"
    icon_p.font.size = Pt(120)
    icon_p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(1.2))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.333), Inches(0.8))
    sub_tf = sub_box.text_frame
    sub_p = sub_tf.paragraphs[0]
    sub_p.text = subtitle
    sub_p.font.size = Pt(28)
    sub_p.font.color.rgb = ACCENT
    sub_p.alignment = PP_ALIGN.CENTER
    
    # Team info
    team_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.333), Inches(0.6))
    team_tf = team_box.text_frame
    team_p = team_tf.paragraphs[0]
    team_p.text = "HackJNU 4.0 | Team Nutshell"
    team_p.font.size = Pt(20)
    team_p.font.color.rgb = RGBColor(180, 180, 180)
    team_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_problem_slide(prs):
    """Slide 2: Problem Statement"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = PRIMARY
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "🔴 The Problem"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Problem points
    problems = [
        ("♿ Accessibility Gap", "Millions with mobility impairments (ALS, cerebral palsy, RSI) cannot use traditional mouse/keyboard"),
        ("💰 Cost Barrier", "Eye-gaze systems cost $10,000+ requiring specialized hardware"),
        ("🔒 Privacy Concerns", "Existing solutions send data to cloud servers"),
        ("🌐 Browser Limitations", "Most assistive tools don't work within web browsers"),
        ("📚 Information Overload", "Clicking countless links just to preview content creates friction")
    ]
    
    y_pos = 1.6
    for icon_title, desc in problems:
        # Point box
        point_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.5), Inches(1))
        tf = point_box.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = icon_title
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(18)
        p2.font.color.rgb = DARK_TEXT
        
        y_pos += 1.1
    
    return slide

def add_solution_slide(prs):
    """Slide 3: Our Solution"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(34, 139, 34)  # Green
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "💡 Our Solution: Nutshell"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(0.5))
    tagline_tf = tagline_box.text_frame
    tagline_p = tagline_tf.paragraphs[0]
    tagline_p.text = "Hands-free browsing powered by head tracking and Chrome Built-in AI"
    tagline_p.font.size = Pt(24)
    tagline_p.font.italic = True
    tagline_p.font.color.rgb = DARK_TEXT
    tagline_p.alignment = PP_ALIGN.CENTER
    
    # Two columns
    # Left column - Head Tracking
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(5.8), Inches(4.5))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "🎯 Head Tracking Control"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    left_features = [
        "• Move cursor with head movements",
        "• Open mouth to click",
        "• Dwell-to-click (hover activation)",
        "• Auto-scroll zones (top/bottom)",
        "• Browser navigation zones"
    ]
    for feat in left_features:
        p = left_tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)
    
    # Right column - AI Summaries
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.4), Inches(5.8), Inches(4.5))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "🤖 AI-Powered Summaries"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    right_features = [
        "• Hover over links for instant previews",
        "• On-device AI (Gemini Nano)",
        "• YouTube video summarization",
        "• Web article key-points",
        "• 100% private - no cloud needed"
    ]
    for feat in right_features:
        p = right_tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)
    
    return slide

def add_features_slide(prs):
    """Slide 4: Key Features"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(70, 130, 180)  # Steel Blue
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "✨ Key Features"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Feature grid (2x3)
    features = [
        ("👄", "Mouth-Open Click", "Calibrated to your facial structure with 800ms cooldown"),
        ("🎯", "Magnetic Snapping", "45px radius auto-targeting of clickable elements"),
        ("📺", "YouTube Support", "Caption extraction & video summarization"),
        ("⚡", "Real-time Streaming", "AI responses stream character by character"),
        ("📷", "Webcam-Based", "Works with any standard camera"),
        ("🔒", "Privacy-First", "100% local processing, zero data sent externally"),
    ]
    
    positions = [
        (0.5, 1.6), (4.6, 1.6), (8.7, 1.6),
        (0.5, 4.4), (4.6, 4.4), (8.7, 4.4)
    ]
    
    for i, (emoji, title, desc) in enumerate(features):
        x, y = positions[i]
        
        # Feature card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RGBColor(200, 200, 200)
        
        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.2), Inches(3.5), Inches(0.6))
        emoji_tf = emoji_box.text_frame
        emoji_p = emoji_tf.paragraphs[0]
        emoji_p.text = emoji
        emoji_p.font.size = Pt(36)
        emoji_p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.85), Inches(3.5), Inches(0.5))
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(18)
        title_p.font.bold = True
        title_p.font.color.rgb = DARK_TEXT
        title_p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.4), Inches(3.3), Inches(1))
        desc_tf = desc_box.text_frame
        desc_tf.word_wrap = True
        desc_p = desc_tf.paragraphs[0]
        desc_p.text = desc
        desc_p.font.size = Pt(14)
        desc_p.font.color.rgb = RGBColor(100, 100, 100)
        desc_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_tech_stack_slide(prs):
    """Slide 5: Technology Stack"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(75, 0, 130)  # Indigo
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "🛠️ Technology Stack"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Tech categories
    categories = [
        ("🎥 Computer Vision", [
            "Human.js - 468-point facial landmark detection",
            "One-Euro Filter - Jitter-free cursor movement",
            "WebGL - GPU-accelerated processing"
        ]),
        ("🤖 AI Integration", [
            "Chrome Summarizer API - Key-point extraction",
            "Chrome Prompt API - Custom prompting",
            "Gemini Nano - On-device language model"
        ]),
        ("🌐 Browser Integration", [
            "Manifest V3 - Modern extension architecture",
            "Side Panel API - Control interface",
            "Readability.js - Content extraction"
        ]),
    ]
    
    x_pos = 0.5
    for cat_title, items in categories:
        cat_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.6), Inches(4), Inches(5))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        
        p = cat_tf.paragraphs[0]
        p.text = cat_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(75, 0, 130)
        
        for item in items:
            p = cat_tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(10)
        
        x_pos += 4.2
    
    # Bottom highlight
    highlight_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(0.8))
    highlight_tf = highlight_box.text_frame
    highlight_p = highlight_tf.paragraphs[0]
    highlight_p.text = "⚡ Lightweight Chrome Extension (~2MB) | Works offline after model download"
    highlight_p.font.size = Pt(20)
    highlight_p.font.color.rgb = ACCENT
    highlight_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_architecture_slide(prs):
    """Slide 6: How It Works (Architecture)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(220, 20, 60)  # Crimson
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "⚙️ How It Works"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Pipeline 1: Head Tracking
    p1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.5), Inches(1.8))
    p1_tf = p1_box.text_frame
    p1_tf.word_wrap = True
    
    p = p1_tf.paragraphs[0]
    p.text = "🎯 Head Tracking Pipeline"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    p = p1_tf.add_paragraph()
    p.text = "Webcam → Face Detection → 468 Landmarks → Head Pose → One-Euro Filter → Screen Coords → Dwell → Action"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_before = Pt(8)
    
    # Pipeline 2: AI Summary
    p2_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.5), Inches(1.8))
    p2_tf = p2_box.text_frame
    p2_tf.word_wrap = True
    
    p = p2_tf.paragraphs[0]
    p.text = "🤖 AI Summary Pipeline"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    p = p2_tf.add_paragraph()
    p.text = "Link Hover (600ms) → Fetch HTML → Readability.js → Chrome AI API → Gemini Nano → Streaming Response → Tooltip"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_before = Pt(8)
    
    # Pipeline 3: YouTube
    p3_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.5), Inches(1.8))
    p3_tf = p3_box.text_frame
    p3_tf.word_wrap = True
    
    p = p3_tf.paragraphs[0]
    p.text = "📺 YouTube Special Pipeline"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    p = p3_tf.add_paragraph()
    p.text = "Page Load → XHR Interceptor → Capture Captions → Parse Text → Combine with Description → Custom Prompt → Summary"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_before = Pt(8)
    
    return slide

def add_innovation_slide(prs):
    """Slide 7: Innovation & Impact"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(255, 140, 0)  # Dark Orange
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "🚀 Innovation & Impact"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Two columns
    # Left: Innovation
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6), Inches(5))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "💡 What Makes It Innovative"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 140, 0)
    
    innovations = [
        "First to combine head tracking + on-device AI",
        "Zero cost ($0) vs. $10,000+ alternatives",
        "100% privacy - all processing is local",
        "Works with any standard webcam",
        "No specialized hardware required",
        "Instant responses - no network latency"
    ]
    for item in innovations:
        p = left_tf.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(10)
    
    # Right: Impact
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(6), Inches(5))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "🌍 Real-World Impact"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 140, 0)
    
    impacts = [
        "Accessibility for mobility impairments",
        "RSI prevention & management",
        "Temporary disability support",
        "Hands-free research & productivity",
        "Democratizes assistive technology",
        "Works offline after initial setup"
    ]
    for item in impacts:
        p = right_tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(10)
    
    return slide

def add_thank_you_slide(prs):
    """Slide 8: Thank You / Contact"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BG
    bg_shape.line.fill.background()
    
    # Emoji
    emoji_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(2.5), Inches(1.2))
    emoji_tf = emoji_box.text_frame
    emoji_p = emoji_tf.paragraphs[0]
    emoji_p.text = "🥜"
    emoji_p.font.size = Pt(80)
    emoji_p.alignment = PP_ALIGN.CENTER
    
    # Thank You
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(11.333), Inches(1))
    thanks_tf = thanks_box.text_frame
    thanks_p = thanks_tf.paragraphs[0]
    thanks_p.text = "Thank You!"
    thanks_p.font.size = Pt(56)
    thanks_p.font.bold = True
    thanks_p.font.color.rgb = WHITE
    thanks_p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.6))
    tagline_tf = tagline_box.text_frame
    tagline_p = tagline_tf.paragraphs[0]
    tagline_p.text = "Browse hands-free, understand faster."
    tagline_p.font.size = Pt(28)
    tagline_p.font.italic = True
    tagline_p.font.color.rgb = ACCENT
    tagline_p.alignment = PP_ALIGN.CENTER
    
    # Demo link
    demo_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.333), Inches(0.6))
    demo_tf = demo_box.text_frame
    demo_p = demo_tf.paragraphs[0]
    demo_p.text = "🎥 Demo: youtu.be/KVOM2VvWypE"
    demo_p.font.size = Pt(22)
    demo_p.font.color.rgb = WHITE
    demo_p.alignment = PP_ALIGN.CENTER
    
    # GitHub link
    github_box = slide.shapes.add_textbox(Inches(1), Inches(5.4), Inches(11.333), Inches(0.6))
    github_tf = github_box.text_frame
    github_p = github_tf.paragraphs[0]
    github_p.text = "📦 GitHub: github.com/priyanshuharshbodhi1/Nutshell"
    github_p.font.size = Pt(22)
    github_p.font.color.rgb = WHITE
    github_p.alignment = PP_ALIGN.CENTER
    
    # HackJNU badge
    badge_box = slide.shapes.add_textbox(Inches(1), Inches(6.4), Inches(11.333), Inches(0.5))
    badge_tf = badge_box.text_frame
    badge_p = badge_tf.paragraphs[0]
    badge_p.text = "HackJNU 4.0 | January 2026"
    badge_p.font.size = Pt(18)
    badge_p.font.color.rgb = RGBColor(150, 150, 150)
    badge_p.alignment = PP_ALIGN.CENTER
    
    return slide


# Create all slides
print("Creating Nutshell Hackathon Presentation...")

add_title_slide(prs, "Nutshell", "Hands-free browsing powered by head tracking & Chrome AI")
add_problem_slide(prs)
add_solution_slide(prs)
add_features_slide(prs)
add_tech_stack_slide(prs)
add_architecture_slide(prs)
add_innovation_slide(prs)
add_thank_you_slide(prs)

# Save presentation
output_path = '/home/priyanshu/repos/Nutshell/Nutshell_HackJNU4_Presentation.pptx'
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
